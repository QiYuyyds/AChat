"""File parsing pipeline — PDF three-level fallback, plain text, DOCX, PPTX.

``parse_bytes()`` remains as the legacy entry point for backward compatibility,
but the module also defines ``DocumentProcessor`` which extends
``BaseDocumentProcessor`` — this is the parser registered in the parser_registry.
"""

from __future__ import annotations

import contextlib
import io
import logging
import os
import shutil
import subprocess
import tempfile

from app.config import get_settings
from app.rag.parsers.base import (
    MIN_USEFUL_PDF_TEXT_RUNES,
    BaseDocumentProcessor,
    ExtractedImage,
    ParseResult,
)

logger = logging.getLogger(__name__)


# ─── Legacy entry point (backward compat) ─────────────────────────────────


def parse_bytes(filename: str, content_type: str, data: bytes) -> ParseResult:
    """Parse uploaded file bytes into a ParseResult.

    Dispatches to PDF parser for .pdf files, DOCX/PPTX parsers for those
    extensions, otherwise decodes as text.
    """
    processor = DocumentProcessor()
    return processor.process_file(filename, content_type, data)


# ─── DocumentProcessor: BaseDocumentProcessor implementation ──────────────


class DocumentProcessor(BaseDocumentProcessor):
    """Default document processor: PDF three-level fallback + text + DOCX/PPTX."""

    supported_extensions = [
        ".txt", ".md", ".markdown", ".csv", ".tsv", ".log",
        ".json", ".xml", ".yaml", ".yml",
        ".pdf", ".docx", ".pptx",
    ]
    service_name = "default"
    display_name = "Default (PDF fallback + text + DOCX/PPTX)"

    def check_health(self) -> bool:
        return True

    def process_file(
        self, filename: str, content_type: str, data: bytes
    ) -> ParseResult:
        filename = filename or ""
        ct = self.normalize_content_type(filename, content_type)
        ext = self.get_extension(filename)

        extract_images = get_settings().rag_extract_images

        if ext == ".pdf":
            return self._parse_pdf(filename, ct, data, extract_images=extract_images)

        if ext == ".docx":
            return self._parse_docx(filename, ct, data, extract_images=extract_images)

        if ext == ".pptx":
            return self._parse_pptx(filename, ct, data, extract_images=extract_images)

        # Fallback: treat as plain text
        return self._parse_text(filename, ct, data)

    # ─── Text parsing ─────────────────────────────────────────────────

    def _parse_text(self, filename: str, ct: str, data: bytes) -> ParseResult:
        text = self.normalize_text(self.decode_text(data))
        if not text.strip():
            raise ValueError("uploaded document is empty")
        return ParseResult(
            filename=filename,
            content_type=ct,
            parser="plain_text",
            content=text,
            text_chars=self.rune_len(text),
        )

    # ─── PDF parsing (three-level fallback) ────────────────────────────

    def _parse_pdf(self, filename: str, ct: str, data: bytes, *, extract_images: bool = True) -> ParseResult:
        """Try pdfplumber → PyPDF2/pypdf → pdftotext in order."""
        for parser_fn in (
            self._extract_pdf_with_pdfplumber,
            self._extract_pdf_with_pypdf2,
            self._extract_pdf_with_pdftotext,
        ):
            text, pages, parser_name, images = parser_fn(data, extract_images=extract_images)
            if text.strip():
                text = self.normalize_text(text)
                chars = self.rune_len(text)
                return ParseResult(
                    filename=filename,
                    content_type=ct,
                    parser=parser_name,
                    content=text,
                    pages=pages,
                    text_chars=chars,
                    needs_ocr=pages > 0 and chars < MIN_USEFUL_PDF_TEXT_RUNES,
                    images=images,
                )
        raise ValueError("pdf contains no extractable text; OCR is required")

    def _extract_pdf_with_pdfplumber(self, data: bytes, *, extract_images: bool = False) -> tuple[str, int, str, list[ExtractedImage]]:
        """Most accurate PDF extractor — requires pdfplumber library."""
        try:
            import pdfplumber  # type: ignore
        except Exception:
            return "", 0, "pdfplumber", []
        try:
            with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as f:
                f.write(data)
                path = f.name
            try:
                texts: list[str] = []
                images: list[ExtractedImage] = []
                with pdfplumber.open(path) as pdf:
                    pages = len(pdf.pages)
                    for idx, page in enumerate(pdf.pages, 1):
                        text = page.extract_text(x_tolerance=1, y_tolerance=3) or ""
                        if text.strip():
                            texts.append(f"--- page {idx} ---\n{text}")
                        if extract_images:
                            for img_idx, img in enumerate(page.images):
                                try:
                                    cropped = page.crop((img["x0"], img["top"], img["x1"], img["bottom"]))
                                    pil_img = cropped.to_image(resolution=150)
                                    buf = io.BytesIO()
                                    pil_img.original.save(buf, format="PNG")
                                    images.append(ExtractedImage(
                                        filename=f"page_{idx}_img_{img_idx + 1}.png",
                                        data=buf.getvalue(),
                                        content_type="image/png",
                                    ))
                                except Exception:
                                    pass
                return "\n\n".join(texts), pages, "pdfplumber", images
            finally:
                _safe_unlink(path)
        except Exception as e:
            logger.warning("pdfplumber extraction failed: %s", e)
            return "", 0, "pdfplumber", []

    def _extract_pdf_with_pypdf2(self, data: bytes, *, extract_images: bool = False) -> tuple[str, int, str, list[ExtractedImage]]:
        """Pure-Python fallback — tries PyPDF2 then pypdf (successor library)."""
        reader = None
        try:
            from PyPDF2 import PdfReader  # type: ignore
            reader = PdfReader(io.BytesIO(data))
        except Exception:
            try:
                from pypdf import PdfReader  # type: ignore
                reader = PdfReader(io.BytesIO(data))
            except Exception:
                return "", 0, "pdf_text", []
        try:
            texts: list[str] = []
            for idx, page in enumerate(reader.pages, 1):
                text = page.extract_text() or ""
                if text.strip():
                    texts.append(f"--- page {idx} ---\n{text}")
            return "\n\n".join(texts), len(reader.pages), "pdf_text", []
        except Exception as e:
            logger.warning("PyPDF2/pypdf extraction failed: %s", e)
            return "", 0, "pdf_text", []

    def _extract_pdf_with_pdftotext(self, data: bytes, *, extract_images: bool = False) -> tuple[str, int, str, list[ExtractedImage]]:
        """System-command fallback — requires pdftotext (poppler-utils) installed."""
        exe = shutil.which("pdftotext")
        if not exe:
            return "", 0, "pdftotext", []
        try:
            with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as f:
                f.write(data)
                path = f.name
            try:
                out = subprocess.check_output(
                    [exe, "-layout", "-enc", "UTF-8", path, "-"],
                    timeout=30,
                )
                text = out.decode("utf-8", errors="ignore")
                pages = text.count("\x0c") + 1 if text.strip() else 0
                return text, pages, "pdftotext", []
            finally:
                _safe_unlink(path)
        except Exception as e:
            logger.warning("pdftotext extraction failed: %s", e)
            return "", 0, "pdftotext", []

    # ─── DOCX parsing ─────────────────────────────────────────────────

    def _parse_docx(self, filename: str, ct: str, data: bytes, *, extract_images: bool = True) -> ParseResult:
        """Extract text from DOCX using python-docx."""
        try:
            from docx import Document as DocxDocument  # type: ignore
        except ImportError:
            logger.warning("python-docx not installed — cannot parse DOCX")
            return ParseResult(
                filename=filename,
                content_type=ct,
                parser="docx",
                needs_ocr=True,
            )

        try:
            doc = DocxDocument(io.BytesIO(data))
            paragraphs: list[str] = []
            for para in doc.paragraphs:
                text = para.text.strip()
                if text:
                    paragraphs.append(text)

            # Also extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_texts: list[str] = []
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        if cell_text:
                            row_texts.append(cell_text)
                    if row_texts:
                        paragraphs.append(" | ".join(row_texts))

            text = self.normalize_text("\n\n".join(paragraphs))

            # Extract inline images
            images: list[ExtractedImage] = []
            if extract_images:
                for idx, shape in enumerate(doc.inline_shapes):
                    try:
                        image = shape._inline.graphic.graphicData.pic.blipFill.blip
                        image_part = doc.part.related_parts[image.embed]
                        ext = image_part.partname.ext
                        content_type_img = "image/png" if ext == ".png" else "image/jpeg"
                        images.append(ExtractedImage(
                            filename=f"image_{idx + 1}{ext}",
                            data=image_part.blob,
                            content_type=content_type_img,
                        ))
                    except Exception:
                        pass

            if not text.strip():
                return ParseResult(
                    filename=filename,
                    content_type=ct,
                    parser="docx",
                    needs_ocr=True,
                    images=images,
                )

            return ParseResult(
                filename=filename,
                content_type=ct,
                parser="docx",
                content=text,
                text_chars=self.rune_len(text),
                images=images,
            )
        except Exception as e:
            logger.warning("DOCX extraction failed for %s: %s", filename, e)
            return ParseResult(
                filename=filename,
                content_type=ct,
                parser="docx",
                needs_ocr=True,
            )

    # ─── PPTX parsing ─────────────────────────────────────────────────

    def _parse_pptx(self, filename: str, ct: str, data: bytes, *, extract_images: bool = True) -> ParseResult:
        """Extract text from PPTX using python-pptx."""
        try:
            from pptx import Presentation  # type: ignore
            from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore
        except ImportError:
            logger.warning("python-pptx not installed — cannot parse PPTX")
            return ParseResult(
                filename=filename,
                content_type=ct,
                parser="pptx",
                needs_ocr=True,
            )

        try:
            prs = Presentation(io.BytesIO(data))
            slides_text: list[str] = []
            images: list[ExtractedImage] = []

            for idx, slide in enumerate(prs.slides, 1):
                texts: list[str] = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        t = shape.text.strip()
                        if t:
                            texts.append(t)
                    # Extract text from tables in slides
                    if shape.has_table:
                        for row in shape.table.rows:
                            row_texts: list[str] = []
                            for cell in row.cells:
                                ct_text = cell.text.strip()
                                if ct_text:
                                    row_texts.append(ct_text)
                            if row_texts:
                                texts.append(" | ".join(row_texts))
                    # Extract images from slides
                    if extract_images and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            image = shape.image
                            images.append(ExtractedImage(
                                filename=f"slide_{idx}_img_{len(images) + 1}.{image.ext}",
                                data=image.blob,
                                content_type=image.content_type,
                            ))
                        except Exception:
                            pass
                if texts:
                    slides_text.append(f"--- slide {idx} ---\n" + "\n".join(texts))

            text = self.normalize_text("\n\n".join(slides_text))

            if not text.strip():
                return ParseResult(
                    filename=filename,
                    content_type=ct,
                    parser="pptx",
                    needs_ocr=True,
                    images=images,
                )

            return ParseResult(
                filename=filename,
                content_type=ct,
                parser="pptx",
                content=text,
                pages=len(prs.slides),
                text_chars=self.rune_len(text),
                images=images,
            )
        except Exception as e:
            logger.warning("PPTX extraction failed for %s: %s", filename, e)
            return ParseResult(
                filename=filename,
                content_type=ct,
                parser="pptx",
                needs_ocr=True,
            )


# ─── Helpers ──────────────────────────────────────────────────────────────


def _safe_unlink(path: str | None) -> None:
    if not path:
        return
    with contextlib.suppress(OSError):
        os.remove(path)
